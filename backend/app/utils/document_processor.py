from fastapi import UploadFile
import PyPDF2
import io
from pptx import Presentation
import os
from pathlib import Path

UPLOAD_DIR = Path("data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

async def process_document(file: UploadFile, document_id: str) -> str:
    """Process uploaded documents and save them"""
    try:
        # Save the original file
        file_path = UPLOAD_DIR / f"{document_id}_{file.filename}"
        content = await file.read()
        
        with open(file_path, "wb") as f:
            f.write(content)
        
        # Extract text based on file type
        text_content = ""
        if file.filename.endswith('.pdf'):
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
                
        elif file.filename.endswith(('.pptx', '.ppt')):
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        
        return text_content, str(file_path)
    
    except Exception as e:
        raise Exception(f"Error processing document: {str(e)}")